"""Generate shootlink-deck.pptx — 10-slide presentation for Shootlink."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Brand Constants ──
PRIMARY = "65195E"
PRIMARY_HOVER = "91157E"
BG = "F8F8F8"
TEXT_DARK = "1A1A1A"
TEXT_MUTED = "666666"
WHITE = "FFFFFF"
ACCENT_GREEN = "059669"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT_TITLE = "Calibri"
FONT_BODY = "Calibri"
FONT_MONO = "Consolas"

OUTPUT_PATH = os.path.join(os.path.dirname(os.path.dirname(__file__)), "docs", "shootlink-deck.pptx")
LOGO_PATH = os.path.join(os.path.dirname(os.path.dirname(__file__)), "public", "logo.png")


def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def set_slide_bg(slide, color_hex):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _resolve_color(color_hex)


def _resolve_color(color):
    """Return RGBColor from either hex string or RGBColor object."""
    if isinstance(color, RGBColor):
        return color
    return hex_to_rgb(color)


def add_shape(slide, left, top, width, height, fill_hex=None, line_hex=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill_hex:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _resolve_color(fill_hex)
    else:
        shape.fill.background()
    if line_hex:
        shape.line.color.rgb = _resolve_color(line_hex)
        shape.line.width = line_width or Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_hex=None, line_hex=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill_hex:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _resolve_color(fill_hex)
    else:
        shape.fill.background()
    if line_hex:
        shape.line.color.rgb = _resolve_color(line_hex)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text="", font_size=14, font_color=TEXT_DARK,
                bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT_BODY, anchor=MSO_ANCHOR.TOP,
                line_spacing=1.15):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.paragraphs[0].alignment = alignment
    except:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = _resolve_color(font_color)
    p.font.bold = bold
    p.font.name = font_name
    p.line_spacing = Pt(int(font_size * line_spacing))
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return txBox


def add_bullet_textbox(slide, left, top, width, height, items, font_size=14, font_color=TEXT_DARK,
                       bold=False, font_name=FONT_BODY, bullet_char="\u2022", spacing_after=4):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet_char}  {item}" if not item.startswith(bullet_char) else item
        p.font.size = Pt(font_size)
        p.font.color.rgb = _resolve_color(font_color)
        p.font.bold = bold
        p.font.name = font_name
        p.space_after = Pt(spacing_after)
        p.line_spacing = Pt(int(font_size * 1.35))
    return txBox


def add_footer(slide, slide_num):
    add_textbox(slide, Inches(0.6), Inches(6.9), Inches(5), Inches(0.4),
                "Shootlink \u2022 Demo", font_size=8, font_color=TEXT_MUTED)
    add_textbox(slide, Inches(11.5), Inches(6.9), Inches(1.5), Inches(0.4),
                str(slide_num), font_size=8, font_color=TEXT_MUTED,
                alignment=PP_ALIGN.RIGHT)


def add_brand_accent_line(slide, top=Inches(0)):
    """Thin primary line at top of content slides."""
    add_shape(slide, Inches(0), top, SLIDE_W, Pt(3), fill_hex=PRIMARY)


def add_section_title(slide, title, subtitle=None):
    add_brand_accent_line(slide)
    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
                title, font_size=26, font_color=PRIMARY, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(1.05), Inches(11), Inches(0.4),
                    subtitle, font_size=13, font_color=TEXT_MUTED)


# ──────────────────────────────────────────────
#  SLIDE BUILDERS
# ──────────────────────────────────────────────

def build_slide_1_cover(prs, n):
    """Cover slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, PRIMARY)

    # Logo
    if os.path.isfile(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(6.1), Inches(1.6), Inches(1.2), Inches(1.2))

    add_textbox(slide, Inches(2), Inches(3.0), Inches(9.3), Inches(0.8),
                "Shootlink", font_size=44, font_color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER, font_name="Calibri")

    tagline = "Galeri Profesional. Alur Kerja Terstruktur. Klien Puas."
    add_textbox(slide, Inches(2), Inches(3.75), Inches(9.3), Inches(0.5),
                tagline, font_size=16, font_color=WHITE,
                alignment=PP_ALIGN.CENTER)

    # Thin separator
    add_shape(slide, Inches(5.5), Inches(4.4), Inches(2.3), Pt(2), fill_hex=WHITE)

    add_textbox(slide, Inches(2), Inches(4.7), Inches(9.3), Inches(0.4),
                "Tim Shootlink", font_size=13, font_color=WHITE,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(2), Inches(5.2), Inches(9.3), Inches(0.4),
                "Platform SaaS Galeri Klien & Proofing Fotografi",
                font_size=11, font_color=RGBColor(0xD0, 0xA0, 0xD0),
                alignment=PP_ALIGN.CENTER)

    add_footer(slide, n)


def build_slide_2_problem_solution(prs, n):
    """Masalah & Solusi — two column layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_section_title(slide, "Masalah & Solusi")

    # LEFT — Masalah
    add_rounded_rect(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(4.8),
                     fill_hex=WHITE, line_hex="E5E5E5")

    add_textbox(slide, Inches(1.1), Inches(1.9), Inches(5), Inches(0.4),
                "Sebelum: Tantangan Fotografer", font_size=15, font_color="CC3333", bold=True)

    left_items = [
        "Review foto campur aduk via WhatsApp & Google Drive",
        "Tidak ada alur kurasi terstruktur",
        "Klien bingung foto mana yang sudah dipilih",
        "Branding lemah — link drive biasa tanpa watermark",
        "Pendapatan hanya dari jasa shooting saja"
    ]
    add_bullet_textbox(slide, Inches(1.1), Inches(2.4), Inches(5), Inches(3.8),
                       left_items, font_size=12, font_color=TEXT_DARK)

    # RIGHT — Solusi
    add_rounded_rect(slide, Inches(6.8), Inches(1.7), Inches(5.8), Inches(4.8),
                     fill_hex=WHITE, line_hex="E5E5E5")

    add_textbox(slide, Inches(7.1), Inches(1.9), Inches(5.3), Inches(0.4),
                "Sesudah: Shootlink", font_size=15, font_color=ACCENT_GREEN, bold=True)

    right_items = [
        "Galeri klien profesional via link unik per proyek",
        "Alur kurasi 6 tahap dari Persiapan hingga Selesai",
        "Pilih favorit, komentar, approve/reject foto langsung",
        "Watermark otomatis proteksi hak cipta",
        "Recurring revenue via subscription bulanan"
    ]
    add_bullet_textbox(slide, Inches(7.1), Inches(2.4), Inches(5.3), Inches(3.8),
                       right_items, font_size=12, font_color=TEXT_DARK)

    add_footer(slide, n)


def build_slide_3_rbac(prs, n):
    """Arsitektur & RBAC."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_section_title(slide, "Arsitektur & RBAC",
                      "Dual-Role Access Control dengan Row Level Security")

    # Role cards
    roles = [
        ("Fotografer (Owner)", PRIMARY, [
            "CRUD proyek & upload foto massal",
            "Atur progress pipeline 6 tahap",
            "Approve / reject hasil kurasi klien",
            "Watermark & download ZIP semua foto",
            "Undang klien via email invite"
        ]),
        ("Klien (Client)", "7C3AED", [
            "Akses galeri via link unik (slug)",
            "Pilih foto favorit tinggalkan komentar",
            "Selesai kurasi — konfirmasi final",
            "Download ZIP foto approved saja",
            "Pantau progress proyek real-time"
        ])
    ]

    for idx, (role_name, color, items) in enumerate(roles):
        x = Inches(0.8 + idx * 6.3)
        add_rounded_rect(slide, x, Inches(1.7), Inches(5.8), Inches(4.5),
                         fill_hex=WHITE, line_hex="E5E5E5")

        # Role badge
        badge = add_rounded_rect(slide, x + Inches(0.3), Inches(2.0), Inches(3.5), Inches(0.4),
                                 fill_hex=color)
        badge.text_frame.paragraphs[0].text = f"  Role: {role_name}"
        badge.text_frame.paragraphs[0].font.size = Pt(11)
        badge.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(WHITE)
        badge.text_frame.paragraphs[0].font.bold = True
        badge.text_frame.paragraphs[0].font.name = FONT_BODY
        badge.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

        add_bullet_textbox(slide, x + Inches(0.3), Inches(2.6), Inches(5.2), Inches(3.2),
                           items, font_size=12, font_color=TEXT_DARK)

    # Tech note at bottom
    add_rounded_rect(slide, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.4),
                     fill_hex=PRIMARY)
    add_textbox(slide, Inches(1.0), Inches(6.4), Inches(11.3), Inches(0.4),
                "RLS (Row Level Security) di Supabase enforce akses per role  |  Hard redirect via window.location.href untuk keamanan routing",
                font_size=9, font_color=WHITE, alignment=PP_ALIGN.CENTER)

    add_footer(slide, n)


def build_slide_4_project_mgmt(prs, n):
    """Manajemen Proyek."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_section_title(slide, "Manajemen Proyek",
                      "Dashboard, Pipeline, dan Organisasi")

    features = [
        ("Dashboard Analitik", "Total proyek, proyek aktif, total foto, foto favorit — update real-time"),
        ("Pipeline 6 Tahap", "Persiapan \u2192 Uploading \u2192 Proses Edit \u2192 Menunggu Review \u2192 Tahap Kurasi Klien \u2192 Selesai"),
        ("Grid & Table View", "Tampilan card atau tabel, filter tab per status, search real-time, sort by name/date/status/progress"),
        ("Buat / Edit / Hapus", "Modal create project dengan jenis acara, edit detail inline, hapus dengan cascade ke storage & relasi")
    ]

    for i, (title, desc) in enumerate(features):
        col = i % 2
        row = i // 2
        x = Inches(0.8 + col * 6.2)
        y = Inches(1.7 + row * 2.3)

        card = add_rounded_rect(slide, x, y, Inches(5.7), Inches(2.0),
                                fill_hex=WHITE, line_hex="E5E5E5")

        # Number badge
        num_shape = add_shape(slide, x + Inches(0.25), y + Inches(0.25), Inches(0.35), Inches(0.35),
                              fill_hex=PRIMARY)
        num_shape.text_frame.paragraphs[0].text = str(i + 1)
        num_shape.text_frame.paragraphs[0].font.size = Pt(12)
        num_shape.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(WHITE)
        num_shape.text_frame.paragraphs[0].font.bold = True
        num_shape.text_frame.paragraphs[0].font.name = FONT_BODY
        num_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_textbox(slide, x + Inches(0.75), y + Inches(0.2), Inches(4.7), Inches(0.35),
                    title, font_size=14, font_color=PRIMARY, bold=True)
        add_textbox(slide, x + Inches(0.75), y + Inches(0.6), Inches(4.7), Inches(1.2),
                    desc, font_size=11, font_color=TEXT_MUTED)

    add_footer(slide, n)


def build_slide_5_curation(prs, n):
    """Galeri & Workflow Kurasi."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_section_title(slide, "Galeri & Workflow Kurasi",
                      "Upload, filter, approve/reject, finalisasi")

    # Left column: Upload & Gallery
    add_rounded_rect(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(4.8),
                     fill_hex=WHITE, line_hex="E5E5E5")

    add_textbox(slide, Inches(1.1), Inches(1.9), Inches(5.3), Inches(0.35),
                "Upload & Galeri", font_size=15, font_color=PRIMARY, bold=True)

    upload_items = [
        "Upload multi-file drag & drop — validasi client-side (JPEG, PNG, WebP, AVIF, max 20MB)",
        "Upload queue sequential dengan status: waiting / uploading / success / error",
        "Token refresh per file — aman untuk upload puluhan foto sekaligus",
        "4 filter tab: Semua / Pilihan Klien / Disetujui / Ditolak",
        "Lightbox detail modal + navigasi prev/next + panel komentar"
    ]
    add_bullet_textbox(slide, Inches(1.1), Inches(2.35), Inches(5.3), Inches(3.8),
                       upload_items, font_size=11, font_color=TEXT_DARK)

    # Right column: Approval Workflow
    add_rounded_rect(slide, Inches(6.8), Inches(1.7), Inches(5.8), Inches(4.8),
                     fill_hex=WHITE, line_hex="E5E5E5")

    add_textbox(slide, Inches(7.1), Inches(1.9), Inches(5.3), Inches(0.35),
                "Approval Workflow", font_size=15, font_color=ACCENT_GREEN, bold=True)

    approval_items = [
        "Setiap foto punya status: pending / approved / rejected",
        "Fotografer approve/reject individual via RPC update_photo_status()",
        "Tombol \"Selesai Kurasi\" — finalisasi massal via RPC finalize_curation()",
        "Semua pending \u2192 approved, progress status \u2192 Selesai",
        "Audit trail lengkap: siapa approve/reject & kapan"
    ]
    add_bullet_textbox(slide, Inches(7.1), Inches(2.35), Inches(5.3), Inches(3.8),
                       approval_items, font_size=11, font_color=TEXT_DARK)

    add_footer(slide, n)


def build_slide_6_collab(prs, n):
    """Fitur Kolaborasi Lanjutan: Comments, Watermark, ZIP."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_section_title(slide, "Fitur Kolaborasi Lanjutan",
                      "Komentar real-time, watermark proteksi, download ZIP batch")

    features = [
        ("Komentar Real-Time", "comments",
         [
             "Thread komentar per foto — real-time via Supabase Realtime",
             "RLS ketat: hanya owner & assigned client yang bisa baca/tulis",
             "Penulis komentar atau owner proyek bisa hapus komentar",
             "Badge jumlah komentar di PhotoCard + panel di detail modal"
         ]),
        ("Watermark Otomatis", PRIMARY,
         [
             "Upload logo watermark sekali — terapkan ke semua foto galeri",
             "9 posisi: top/center/bottom + left/center/right",
             "Slider opasitas (10%-100%) & ukuran (5%-50%)",
             "Overlay client-side via Canvas — tidak merusak file original"
         ]),
        ("Download ZIP Batch", "7C3AED",
         [
             "Owner: download semua foto dalam 1 file ZIP",
             "Client: download hanya foto berstatus \"approved\"",
             "Streaming via archiver — memory efficient, tanpa file temp",
             "Nama file: {nama-proyek}-foto.zip"
         ])
    ]

    for i, (title, color, items) in enumerate(features):
        x = Inches(0.8 + i * 4.15)
        add_rounded_rect(slide, x, Inches(1.7), Inches(3.85), Inches(4.5),
                         fill_hex=WHITE, line_hex="E5E5E5")

        # Title badge
        badge_color = color if color and color != "comments" else "3B82F6"
        badge = add_rounded_rect(slide, x + Inches(0.25), Inches(1.95), Inches(3.35), Inches(0.4),
                                 fill_hex=badge_color)
        badge.text_frame.paragraphs[0].text = f"  {title}"
        badge.text_frame.paragraphs[0].font.size = Pt(11)
        badge.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(WHITE)
        badge.text_frame.paragraphs[0].font.bold = True
        badge.text_frame.paragraphs[0].font.name = FONT_BODY
        badge.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

        add_bullet_textbox(slide, x + Inches(0.25), Inches(2.55), Inches(3.35), Inches(3.2),
                           items, font_size=11, font_color=TEXT_DARK)

    add_footer(slide, n)


def build_slide_7_client_experience(prs, n):
    """Client Experience & Invite Flow."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_section_title(slide, "Client Experience & Invite",
                      "Dari undangan hingga finalisasi — semua dalam 1 link")

    # Left: Invite Flow
    add_rounded_rect(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(4.8),
                     fill_hex=WHITE, line_hex="E5E5E5")

    add_textbox(slide, Inches(1.1), Inches(1.9), Inches(5.3), Inches(0.35),
                "Flow Undangan Klien", font_size=15, font_color=PRIMARY, bold=True)

    flow_steps = [
        "1. Fotografer klik \"Undang Klien\" di toolbar galeri",
        "2. Masukkan email klien \u2192 insert ke tabel project_clients",
        "3. Klien terima email invite dengan link akses",
        "4. Klien signup/login \u2192 accepted_at terisi \u2192 akses granted",
        "5. RLS enforce: hanya assigned client bisa lihat proyek itu"
    ]
    add_bullet_textbox(slide, Inches(1.1), Inches(2.35), Inches(5.3), Inches(3.8),
                       flow_steps, font_size=12, font_color=TEXT_DARK)

    # Right: Client Gallery
    add_rounded_rect(slide, Inches(6.8), Inches(1.7), Inches(5.8), Inches(4.8),
                     fill_hex=WHITE, line_hex="E5E5E5")

    add_textbox(slide, Inches(7.1), Inches(1.9), Inches(5.3), Inches(0.35),
                "Pengalaman Client", font_size=15, font_color="7C3AED", bold=True)

    client_items = [
        "Akses galeri via URL: domain.com/{unique-slug}",
        "Favoritkan foto yang disukai (toggle real-time)",
        "Tinggalkan komentar per foto untuk revisi",
        "Filter: Semua / Pilihan Saya / Disetujui / Ditolak",
        "Klik \"Selesai Kurasi\" \u2192 finalisasi pilihan",
        "Download ZIP foto yang sudah disetujui"
    ]
    add_bullet_textbox(slide, Inches(7.1), Inches(2.35), Inches(5.3), Inches(3.8),
                       client_items, font_size=12, font_color=TEXT_DARK)

    add_footer(slide, n)


def build_slide_8_pricing(prs, n):
    """Monetisasi & Subscription."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_section_title(slide, "Monetisasi & Subscription",
                      "3 tier pricing dengan Stripe Checkout + plan limits otomatis")

    plans = [
        ("Free", "Gratis", "", PRIMARY, [
            "3 proyek aktif",
            "100 foto per proyek",
            "Gallery client standar",
            "Progress tracking",
            "Bagikan tautan galeri"
        ], False),
        ("Basic", "Rp199.000", "/bulan", PRIMARY, [
            "20 proyek aktif",
            "500 foto per proyek",
            "Gallery client + kurasi",
            "Undang klien via email",
            "Prioritas support"
        ], True),
        ("Pro", "Rp499.000", "/bulan", PRIMARY, [
            "Proyek tak terbatas",
            "Foto tak terbatas",
            "Gallery client + kurasi",
            "Custom branding",
            "Prioritas support 24/7"
        ], False)
    ]

    for i, (name, price, period, color, features, popular) in enumerate(plans):
        x = Inches(0.8 + i * 4.15)
        y = Inches(1.6)

        card = add_rounded_rect(slide, x, y, Inches(3.85), Inches(5.0),
                                fill_hex=WHITE, line_hex=color if popular else "E5E5E5")

        if popular:
            # Popular badge
            pop_badge = add_rounded_rect(slide, x + Inches(0.6), y - Inches(0.15),
                                         Inches(2.65), Inches(0.35), fill_hex=color)
            pop_badge.text_frame.paragraphs[0].text = "  Terpopuler"
            pop_badge.text_frame.paragraphs[0].font.size = Pt(9)
            pop_badge.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(WHITE)
            pop_badge.text_frame.paragraphs[0].font.bold = True
            pop_badge.text_frame.paragraphs[0].font.name = FONT_BODY
            pop_badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            # Raise card top to accommodate badge
            y2 = Inches(1.9)
        else:
            y2 = y + Inches(0.3)

        add_textbox(slide, x + Inches(0.3), y2, Inches(3.25), Inches(0.4),
                    name, font_size=18, font_color=color, bold=True,
                    alignment=PP_ALIGN.CENTER)

        add_textbox(slide, x + Inches(0.3), y2 + Inches(0.45), Inches(3.25), Inches(0.4),
                    price, font_size=22, font_color=TEXT_DARK, bold=True,
                    alignment=PP_ALIGN.CENTER)

        if period:
            add_textbox(slide, x + Inches(0.3), y2 + Inches(0.85), Inches(3.25), Inches(0.3),
                        period, font_size=11, font_color=TEXT_MUTED,
                        alignment=PP_ALIGN.CENTER)

        # Features
        feat_y = y2 + Inches(1.3)
        for fi, feat in enumerate(features):
            add_textbox(slide, x + Inches(0.4), Inches(feat_y), Inches(3.1), Inches(0.3),
                        f"\u2713  {feat}", font_size=10, font_color=TEXT_DARK)

    # Tech note
    add_rounded_rect(slide, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.35),
                     fill_hex=PRIMARY)
    add_textbox(slide, Inches(1.0), Inches(6.7), Inches(11.3), Inches(0.35),
                "Stripe Checkout Session  |  Webhook update subscription  |  Plan limits via DB function check_project_limit()",
                font_size=9, font_color=WHITE, alignment=PP_ALIGN.CENTER)

    add_footer(slide, n)


def build_slide_9_tech(prs, n):
    """Teknologi & Kualitas."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_section_title(slide, "Teknologi & Kualitas",
                      "Stack modern, testing ketat, siap production")

    # Left: Tech Stack
    add_rounded_rect(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(4.5),
                     fill_hex=WHITE, line_hex="E5E5E5")

    add_textbox(slide, Inches(1.1), Inches(1.9), Inches(5.3), Inches(0.35),
                "Tech Stack", font_size=15, font_color=PRIMARY, bold=True)

    tech_items = [
        "Frontend: Next.js 15 (App Router) + React 19 + TypeScript (strict)",
        "Styling: Tailwind CSS + RemixIcon (CDN icons)",
        "Database: Supabase PostgreSQL + RLS (Row Level Security)",
        "Auth: Supabase Auth (Email/Password + Magic Link)",
        "Storage: Supabase Storage bucket photos + avatars",
        "Realtime: Supabase Realtime untuk update progress & komentar",
        "Billing: Stripe Checkout Session + Webhook integration"
    ]
    add_bullet_textbox(slide, Inches(1.1), Inches(2.35), Inches(5.3), Inches(3.5),
                       tech_items, font_size=11, font_color=TEXT_DARK)

    # Right: Testing
    add_rounded_rect(slide, Inches(6.8), Inches(1.7), Inches(5.8), Inches(4.5),
                     fill_hex=WHITE, line_hex="E5E5E5")

    add_textbox(slide, Inches(7.1), Inches(1.9), Inches(5.3), Inches(0.35),
                "Testing & Quality", font_size=15, font_color=ACCENT_GREEN, bold=True)

    test_items = [
        "Playwright E2E: 62 test cases (Chromium + Mobile Chrome)",
        "Pass rate: 96.8% (60/62 passing)",
        "Cakupan: Auth, RBAC, Project CRUD, Gallery, Upload, Approval, Comments, Watermark, ZIP, Delete, Invite",
        "Build production: npm run build sukses",
        "Keamanan: function search_path fixed, public execute revoked, storage bucket listing blocked"
    ]
    add_bullet_textbox(slide, Inches(7.1), Inches(2.35), Inches(5.3), Inches(3.5),
                       test_items, font_size=11, font_color=TEXT_DARK)

    add_footer(slide, n)


def build_slide_10_closing(prs, n):
    """Roadmap & Penutup."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PRIMARY)

    add_textbox(slide, Inches(2), Inches(1.2), Inches(9.3), Inches(0.7),
                "Roadmap & Penutup", font_size=32, font_color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)

    # Thin separator
    add_shape(slide, Inches(5.5), Inches(2.0), Inches(2.3), Pt(2), fill_hex=WHITE)

    # Roadmap
    roadmap_items = [
        "v1.1  Email notifikasi & custom domain",
        "v1.2  Before/After slider & AI-assisted culling",
        "v2.0  Mobile app & marketplace fotografer"
    ]
    add_bullet_textbox(slide, Inches(3.5), Inches(2.3), Inches(6.3), Inches(1.8),
                       roadmap_items, font_size=16, font_color=WHITE,
                       bullet_char="\u25B8", spacing_after=8)

    add_shape(slide, Inches(4), Inches(4.2), Inches(5.3), Pt(1), fill_hex=RGBColor(0xD0, 0xA0, 0xD0))

    add_textbox(slide, Inches(2), Inches(4.5), Inches(9.3), Inches(0.5),
                "Shootlink bukan cuma galeri foto \u2014 ini sistem operasi bisnis fotografi.",
                font_size=15, font_color=WHITE, alignment=PP_ALIGN.CENTER)

    # CTA
    cta = add_rounded_rect(slide, Inches(4.5), Inches(5.3), Inches(4.3), Inches(0.55),
                            fill_hex=WHITE)
    cta.text_frame.paragraphs[0].text = "  Siap Demo Langsung?  \u2192"
    cta.text_frame.paragraphs[0].font.size = Pt(14)
    cta.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(PRIMARY)
    cta.text_frame.paragraphs[0].font.bold = True
    cta.text_frame.paragraphs[0].font.name = FONT_BODY
    cta.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_textbox(slide, Inches(2), Inches(6.2), Inches(9.3), Inches(0.4),
                "Terima Kasih \u2022 Ada Pertanyaan?",
                font_size=14, font_color=RGBColor(0xD0, 0xA0, 0xD0),
                alignment=PP_ALIGN.CENTER)

    add_footer(slide, n)


# ── Main ──
def generate():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        build_slide_1_cover,
        build_slide_2_problem_solution,
        build_slide_3_rbac,
        build_slide_4_project_mgmt,
        build_slide_5_curation,
        build_slide_6_collab,
        build_slide_7_client_experience,
        build_slide_8_pricing,
        build_slide_9_tech,
        build_slide_10_closing,
    ]

    for i, builder in enumerate(builders, 1):
        builder(prs, i)

    prs.save(OUTPUT_PATH)
    print(f"Saved: {OUTPUT_PATH}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    generate()
